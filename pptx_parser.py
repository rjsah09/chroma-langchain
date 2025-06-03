
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

'''
작성자: 양건모
최초 생성일: 2025.06.02
-----------------------
변경 내역
2025.06.02 주석 추가
'''
class PPTXParser:
    def __init__(self, pptx_path: str):
        '''변수 및 관련 모듈 초기화'''
        self.pptx_path = pptx_path
        self.pptx = Presentation(pptx_path)
        self.file_name = Path(pptx_path).name

    def extract_shape(self, shape):
        '''shape에서 타입에 맞는 텍스트 추출'''
        if shape.has_text_frame:
            text = shape.text.strip()
            if not text:
                return None
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "text",
                "content": shape.text.strip()
            }
        elif shape.has_table:
            texts = [
                cell.text.strip()
                for row in shape.table.rows
                for cell in row.cells
                if cell.text.strip()
            ]
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "text",
                "content": "\n".join(texts)
            }
        elif shape.has_chart and shape.chart.chart_title:
            title = shape.chart.chart_title.text_frame.text.strip()
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "text",
                "content": title
            }
        elif hasattr(shape, "smartArt"):
            texts = [
                node.text_frame.text.strip()
                for node in shape.smartArt.nodes
                if node.text_frame.text.strip()
            ]
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "text",
                "content": "\n".join(texts)
            }
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_items = []
            for sub_shape in shape.shapes:
                res = self.extract_shape(sub_shape)
                if res:
                    group_items.append(res)
            group_items.sort(key=lambda x: (x["top"], x["left"]))
            return {
                "top": shape.top,
                "left": shape.left,
                "content_type": "group",
                "content": group_items
            }
        return None

    def parse_slides(self):
        '''파일에서 슬라이드 단위 텍스트 파싱'''
        all_slides = []
        for slide_idx, slide in enumerate(self.pptx.slides, start=1):
            items = []
            for shape in slide.shapes:
                extract_result = self.extract_shape(shape)
                if extract_result:
                    items.append(extract_result)
            items.sort(key=lambda x: (x["top"], x["left"]))
            slide_info = {
                "slide_number": slide_idx,
                "file_name": self.file_name,
                "content_items": items
            }
            all_slides.append(slide_info)
        return all_slides
